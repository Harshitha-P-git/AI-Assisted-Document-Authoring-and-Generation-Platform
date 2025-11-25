from fastapi import APIRouter, Depends, HTTPException, status
from fastapi.responses import FileResponse, StreamingResponse
from sqlalchemy.orm import Session
from app.db.database import get_db
from app.db import models
from app.api.v1.endpoints.auth import get_current_user
from app.services.storage_service import save_file
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import logging
from pathlib import Path

logger = logging.getLogger(__name__)
router = APIRouter()


def verify_project_ownership(project_id: int, user_id: int, db: Session) -> models.Project:
    """Verify project ownership"""
    project = db.query(models.Project).filter(
        models.Project.id == project_id,
        models.Project.owner_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Project not found"
        )
    
    return project


@router.get("/project/{project_id}/word")
async def export_word_document(
    project_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Export Word document (.docx)"""
    project = verify_project_ownership(project_id, current_user.id, db)
    
    if project.project_type != "word":
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Project is not a Word project"
        )
    
    # Get sections
    sections = db.query(models.Section).filter(
        models.Section.project_id == project_id
    ).order_by(models.Section.order_index).all()
    
    if not sections:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No sections found in project"
        )
    
    # Create Word document
    doc = Document()
    
    # Add title
    title = doc.add_heading(project.name, 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    if project.description:
        desc = doc.add_paragraph(project.description)
        desc.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph()  # Add spacing
    
    # Add sections
    for section in sections:
        # Add section title
        heading = doc.add_heading(section.title, level=1)
        
        # Add section content
        if section.content:
            # Split content into paragraphs
            paragraphs = section.content.split('\n\n')
            for para_text in paragraphs:
                if para_text.strip():
                    para = doc.add_paragraph(para_text.strip())
                    para_format = para.paragraph_format
                    para_format.space_after = Pt(12)
        
        doc.add_paragraph()  # Add spacing between sections
    
    # Save to BytesIO
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    
    # Save to storage
    filename = f"{project.name.replace(' ', '_')}_{project_id}.docx"
    file_path = f"exports/{current_user.id}/{filename}"
    
    try:
        save_file(file_path, buffer.getvalue())
    except Exception as e:
        logger.error(f"Error saving file: {e}")
    
    # Return file with explicit CORS headers
    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={
            "Content-Disposition": f'attachment; filename="{filename}"',
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Methods": "GET, OPTIONS",
            "Access-Control-Allow-Headers": "*",
        }
    )


@router.get("/project/{project_id}/powerpoint")
async def export_powerpoint_document(
    project_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Export PowerPoint document (.pptx)"""
    try:
        project = verify_project_ownership(project_id, current_user.id, db)
        
        if project.project_type != "powerpoint":
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Project is not a PowerPoint project"
            )
        
        # Get slides
        slides = db.query(models.Slide).filter(
            models.Slide.project_id == project_id
        ).order_by(models.Slide.order_index).all()
        
        if not slides:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="No slides found in project"
            )
        
        logger.info(f"Exporting PowerPoint for project {project_id} with {len(slides)} slides")
        
        # Create PowerPoint presentation with better dimensions
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        try:
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = project.name
            if project.description:
                subtitle.text = project.description
        except Exception as e:
            logger.warning(f"Error creating title slide: {e}")
            # Continue even if title slide fails
        
        # Add content slides
        for slide_data in slides:
            try:
                # Use blank layout
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Add title
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                height = Inches(1)
                
                title_box = slide.shapes.add_textbox(left, top, width, height)
                title_frame = title_box.text_frame
                title_frame.text = slide_data.title or "Untitled Slide"
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(36)
                title_para.font.bold = True
                title_para.font.color.rgb = None  # Use theme color
                title_para.space_after = Pt(12)
                
                # Add content
                if slide_data.content:
                    content_left = Inches(0.5)
                    content_top = Inches(2)
                    content_width = Inches(9)
                    content_height = Inches(5)
                    
                    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                    content_frame = content_box.text_frame
                    content_frame.word_wrap = True
                    
                    # Split content into paragraphs/bullets with better formatting
                    lines = slide_data.content.split('\n')
                    for i, line in enumerate(lines):
                        line = line.strip()
                        if line:
                            # Skip disclaimer lines
                            if "fallback" in line.lower() or "python 3.14" in line.lower() or "compatibility" in line.lower():
                                continue
                            if i == 0:
                                p = content_frame.paragraphs[0]
                            else:
                                p = content_frame.add_paragraph()
                            # Format bullet points
                            if line.startswith('•') or line.startswith('-'):
                                p.text = line
                                p.level = 0
                            else:
                                p.text = line
                                p.level = 0
                            p.font.size = Pt(20)
                            p.space_after = Pt(8)
                            # Make first line slightly larger
                            if i == 0:
                                p.font.size = Pt(22)
            except Exception as e:
                logger.error(f"Error creating slide {slide_data.id}: {e}")
                # Continue with next slide even if one fails
        
        # Save to BytesIO
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        # Save to storage
        filename = f"{project.name.replace(' ', '_')}_{project_id}.pptx"
        file_path = f"exports/{current_user.id}/{filename}"
        
        try:
            save_file(file_path, buffer.getvalue())
        except Exception as e:
            logger.error(f"Error saving file: {e}")
        
        # Return file with explicit CORS headers
        buffer.seek(0)
        return StreamingResponse(
            buffer,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}"',
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Methods": "GET, OPTIONS",
                "Access-Control-Allow-Headers": "*",
            }
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error exporting PowerPoint: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to export PowerPoint: {str(e)}"
        )

